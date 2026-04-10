"""Upload / Paste router for ingesting documents into the knowledge store."""

from __future__ import annotations

import io
import logging
import os
import uuid
from typing import Any, List, Optional

from fastapi import APIRouter, File, Form, HTTPException, Request, UploadFile
from fastapi.responses import StreamingResponse
from pydantic import BaseModel

from openjarvis.connectors.store import KnowledgeStore
from openjarvis.core.config import DEFAULT_CONFIG_DIR
from openjarvis.server.auth import (
    get_operator_memory_manager,
    require_current_user_if_bootstrapped,
)

logger = logging.getLogger(__name__)

router = APIRouter(prefix="/v1/connectors/upload", tags=["upload"])

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

_ALLOWED_EXTENSIONS = {".txt", ".md", ".csv", ".tsv", ".pdf", ".docx", ".xlsx", ".pptx"}


def _chunk_text(text: str, max_chars: int = 1000) -> List[str]:
    """Split *text* into ~max_chars pieces at paragraph boundaries."""
    paragraphs = text.split("\n\n")
    chunks: List[str] = []
    current = ""
    for para in paragraphs:
        para = para.strip()
        if not para:
            continue
        if current and len(current) + len(para) + 2 > max_chars:
            chunks.append(current.strip())
            current = para
        else:
            current = f"{current}\n\n{para}" if current else para
    if current.strip():
        chunks.append(current.strip())
    # Guard against very large paragraphs that exceed max_chars
    final: List[str] = []
    for chunk in chunks:
        while len(chunk) > max_chars:
            # Find last space within limit
            split_at = chunk.rfind(" ", 0, max_chars)
            if split_at == -1:
                split_at = max_chars
            final.append(chunk[:split_at].strip())
            chunk = chunk[split_at:].strip()
        if chunk:
            final.append(chunk)
    return final


def _extract_text_from_pdf(data: bytes) -> str:
    """Extract text from a PDF using pdfplumber or PyPDF2."""
    # Try pdfplumber first
    try:
        import pdfplumber  # type: ignore[import-untyped]

        with pdfplumber.open(io.BytesIO(data)) as pdf:
            pages = [p.extract_text() or "" for p in pdf.pages]
        return "\n\n".join(pages)
    except ImportError:
        pass

    # Fall back to PyPDF2
    try:
        from PyPDF2 import PdfReader  # type: ignore[import-untyped]

        reader = PdfReader(io.BytesIO(data))
        pages = [p.extract_text() or "" for p in reader.pages]
        return "\n\n".join(pages)
    except ImportError:
        raise HTTPException(
            status_code=500,
            detail=(
                "PDF parsing requires pdfplumber or PyPDF2. "
                "Install one with: pip install pdfplumber"
            ),
        )


def _extract_text_from_docx(data: bytes) -> str:
    """Extract text from a .docx file using python-docx."""
    try:
        from docx import Document  # type: ignore[import-untyped]

        doc = Document(io.BytesIO(data))
        return "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())
    except ImportError:
        raise HTTPException(
            status_code=500,
            detail=(
                "DOCX parsing requires python-docx. "
                "Install with: pip install python-docx"
            ),
        )


def _extract_text_from_xlsx(data: bytes) -> str:
    """Extract readable text from an .xlsx workbook using openpyxl."""
    try:
        from openpyxl import load_workbook  # type: ignore[import-untyped]

        workbook = load_workbook(io.BytesIO(data), data_only=True, read_only=True)
        sections: list[str] = []
        for sheet in workbook.worksheets:
            rows: list[str] = []
            for row in sheet.iter_rows(values_only=True):
                cells = [str(value).strip() for value in row if value is not None and str(value).strip()]
                if cells:
                    rows.append(" | ".join(cells))
            if rows:
                sections.append(f"Sheet: {sheet.title}\n" + "\n".join(rows))
        return "\n\n".join(sections)
    except ImportError:
        raise HTTPException(
            status_code=500,
            detail=(
                "XLSX parsing requires openpyxl. "
                "Install with: pip install openpyxl"
            ),
        )


def _extract_text_from_pptx(data: bytes) -> str:
    """Extract readable text from a .pptx deck using python-pptx."""
    try:
        from pptx import Presentation  # type: ignore[import-untyped]

        presentation = Presentation(io.BytesIO(data))
        slides: list[str] = []
        for index, slide in enumerate(presentation.slides, start=1):
            texts: list[str] = []
            for shape in slide.shapes:
                text = getattr(shape, "text", "")
                if text and str(text).strip():
                    texts.append(str(text).strip())
            if texts:
                slides.append(f"Slide {index}\n" + "\n".join(texts))
        return "\n\n".join(slides)
    except ImportError:
        raise HTTPException(
            status_code=500,
            detail=(
                "PPTX parsing requires python-pptx. "
                "Install with: pip install python-pptx"
            ),
        )


def _get_store() -> KnowledgeStore:
    """Return a KnowledgeStore pointing at the default knowledge DB."""
    db_path = DEFAULT_CONFIG_DIR / "knowledge.db"
    return KnowledgeStore(db_path=db_path)


def _owner_user_id(request: Request) -> str:
    user = require_current_user_if_bootstrapped(request)
    if user is None:
        return ""
    return str(user.get("id", "")).strip()


# ---------------------------------------------------------------------------
# Schemas
# ---------------------------------------------------------------------------


class PasteRequest(BaseModel):
    title: str = ""
    content: str


class IngestResponse(BaseModel):
    chunks_added: int
    source: str = "upload"


class DocumentAnalysisResponse(BaseModel):
    mode: str
    content: str
    files: list[str]
    model: str


class DocumentExportRequest(BaseModel):
    title: str = "Document Intel Export"
    mode: str = "summary"
    content: str
    format: str = "docx"


def _update_document_mission(
    request: Request,
    *,
    title: str,
    summary: str,
    next_step: str,
    result: str = "",
    status: str = "active",
    phase: str = "plan",
    retry_hint: str = "",
    result_data: dict[str, Any] | None = None,
    next_action: dict[str, Any] | None = None,
) -> None:
    operator_memory = get_operator_memory_manager(request)
    mission_result_data = {
        "summary": summary,
        "result": result,
        "phase": phase,
        "status": status,
        **(result_data or {}),
    }
    mode = str(mission_result_data.get("mode", "")).strip().lower()
    default_kind = "brief" if mode in {"business_review", "finance_review", "investment_memo", "kpi_extract"} else "prompt"
    default_label = (
        f"{title} Memo"
        if default_kind == "brief" and status == "complete"
        else title
    )
    try:
        operator_memory.update_mission(
            "document-mission",
            {
                "title": title,
                "domain": "document",
                "status": status,
                "phase": phase,
                "summary": summary,
                "next_step": next_step,
                "result": result,
                "retry_hint": retry_hint,
                "result_data": mission_result_data,
                "next_action": next_action
                or {
                    "kind": default_kind,
                    "content": result or next_step or summary,
                    "label": default_label,
                },
            },
        )
    except Exception:
        logger.debug("Document mission update skipped", exc_info=True)


def _safe_export_stem(title: str) -> str:
    cleaned = "".join(ch if ch.isalnum() or ch in {"-", "_", " "} else "-" for ch in title).strip()
    cleaned = "-".join(cleaned.split())
    return cleaned or "document-intel-export"


def _export_document_docx(title: str, mode: str, content: str) -> bytes:
    try:
        from docx import Document  # type: ignore[import-untyped]
    except ImportError as exc:
        raise HTTPException(
            status_code=500,
            detail="DOCX export requires python-docx. Install with: pip install python-docx",
        ) from exc
    doc = Document()
    doc.add_heading(title, level=1)
    doc.add_paragraph(f"Mode: {mode.replace('_', ' ')}")
    for block in [part.strip() for part in content.split("\n\n") if part.strip()]:
        doc.add_paragraph(block)
    buffer = io.BytesIO()
    doc.save(buffer)
    return buffer.getvalue()


def _export_document_txt(title: str, mode: str, content: str) -> bytes:
    payload = f"{title}\nMode: {mode.replace('_', ' ')}\n\n{content}".strip() + "\n"
    return payload.encode("utf-8")


def _export_document_xlsx(title: str, mode: str, content: str) -> bytes:
    try:
        from openpyxl import Workbook  # type: ignore[import-untyped]
    except ImportError as exc:
        raise HTTPException(
            status_code=500,
            detail="XLSX export requires openpyxl. Install with: pip install openpyxl",
        ) from exc
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Document Intel"
    sheet.append(["Title", title])
    sheet.append(["Mode", mode.replace("_", " ")])
    sheet.append([])
    sheet.append(["Section", "Detail"])
    current_section = "Notes"
    for line in content.splitlines():
        cleaned = line.strip()
        if not cleaned:
            continue
        if cleaned.endswith(":") and len(cleaned) < 80:
            current_section = cleaned[:-1]
            continue
        if cleaned.startswith(("-", "*", "•")):
            cleaned = cleaned.lstrip("-*• ").strip()
        sheet.append([current_section, cleaned])
    buffer = io.BytesIO()
    workbook.save(buffer)
    return buffer.getvalue()


# ---------------------------------------------------------------------------
# Routes
# ---------------------------------------------------------------------------


@router.post("/ingest", response_model=IngestResponse)
async def ingest_paste(body: PasteRequest, request: Request) -> IngestResponse:
    """Ingest pasted text into the knowledge store."""
    text = body.content.strip()
    if not text:
        raise HTTPException(status_code=400, detail="Content is empty")

    store = _get_store()
    doc_id = str(uuid.uuid4())
    chunks = _chunk_text(text)
    owner_user_id = _owner_user_id(request)

    for idx, chunk in enumerate(chunks):
        store.store(
            chunk,
            source="upload",
            doc_type="paste",
            doc_id=doc_id,
            title=body.title or "Pasted text",
            chunk_index=idx,
            owner_user_id=owner_user_id,
        )

    logger.info("Ingested %d chunks from pasted text (doc_id=%s)", len(chunks), doc_id)
    return IngestResponse(chunks_added=len(chunks))


@router.post("/ingest/files", response_model=IngestResponse)
async def ingest_files(
    request: Request,
    files: List[UploadFile] = File(...),
    title: Optional[str] = Form(None),
) -> IngestResponse:
    """Ingest uploaded files into the knowledge store."""
    store = _get_store()
    total_chunks = 0
    owner_user_id = _owner_user_id(request)

    for upload in files:
        filename = upload.filename or "untitled"
        ext = ""
        if "." in filename:
            ext = "." + filename.rsplit(".", 1)[-1].lower()

        if ext not in _ALLOWED_EXTENSIONS:
            allowed = ", ".join(sorted(_ALLOWED_EXTENSIONS))
            raise HTTPException(
                status_code=400,
                detail=(f"Unsupported file type: {ext}. Allowed: {allowed}"),
            )

        data = await upload.read()

        text = _extract_text_for_extension(ext, data)

        text = text.strip()
        if not text:
            continue

        doc_id = str(uuid.uuid4())
        doc_title = title or filename
        chunks = _chunk_text(text)

        for idx, chunk in enumerate(chunks):
            store.store(
                chunk,
                source="upload",
                doc_type=ext.lstrip("."),
                doc_id=doc_id,
                title=doc_title,
                chunk_index=idx,
                owner_user_id=owner_user_id,
            )

        total_chunks += len(chunks)
        logger.info(
            "Ingested %d chunks from file %s (doc_id=%s)",
            len(chunks),
            filename,
            doc_id,
        )

    return IngestResponse(chunks_added=total_chunks)


def _extract_text_for_extension(ext: str, data: bytes) -> str:
    if ext in (".txt", ".md", ".csv", ".tsv"):
        try:
            return data.decode("utf-8")
        except UnicodeDecodeError:
            return data.decode("latin-1")
    if ext == ".pdf":
        return _extract_text_from_pdf(data)
    if ext == ".docx":
        return _extract_text_from_docx(data)
    if ext == ".xlsx":
        return _extract_text_from_xlsx(data)
    if ext == ".pptx":
        return _extract_text_from_pptx(data)
    return ""


def _document_analysis_prompt(mode: str, title: str, filenames: list[str]) -> str:
    label = title.strip() or ", ".join(filenames[:3]) or "document set"
    prompts = {
        "summary": (
            "You are JARVIS document intelligence. Reply in English only. "
            "Summarize the provided documents clearly for an operator. "
            "Use sections: Executive Summary, Key Points, Risks, Recommended Next Step."
        ),
        "business_review": (
            "You are JARVIS business analyst mode. Reply in English only. "
            "Review the provided documents like a strong business operator. "
            "Use sections: Executive Summary, Business Model, Operating Signals, Risks, Open Questions, Recommended Next Step."
        ),
        "finance_review": (
            "You are JARVIS finance analyst mode. Reply in English only. "
            "Review the provided documents for financial meaning. "
            "Extract concrete numbers when visible, distinguish facts from assumptions, and use sections: "
            "Executive Summary, Key Metrics, Trends and Deltas, Risks, Missing Data, Recommended Next Step."
        ),
        "investment_memo": (
            "You are JARVIS investment analyst mode. Reply in English only. "
            "Turn the provided materials into an investor-style review. "
            "Use sections: Investment Thesis, Key Metrics, Strengths, Risks, Open Questions, Recommendation."
        ),
        "kpi_extract": (
            "You are JARVIS financial extraction mode. Reply in English only. "
            "Extract concrete KPIs, metrics, tabular signals, trends, and deltas from the provided materials. "
            "Prefer numbers and units exactly as visible. "
            "Use sections: KPI Snapshot, Trends and Deltas, Missing Metrics, Risks or Anomalies, Recommended Next Step."
        ),
    }
    prompt = prompts.get(mode, prompts["summary"])
    return f"{prompt}\n\nDocument label: {label}\nFiles: {', '.join(filenames) or 'unknown'}"


@router.post("/analyze/files", response_model=DocumentAnalysisResponse)
async def analyze_files(
    request: Request,
    files: List[UploadFile] = File(...),
    mode: str = Form("summary"),
    title: Optional[str] = Form(None),
) -> DocumentAnalysisResponse:
    allowed_modes = {"summary", "business_review", "finance_review", "investment_memo", "kpi_extract"}
    selected_mode = mode.strip().lower()
    if selected_mode not in allowed_modes:
        raise HTTPException(status_code=400, detail=f"Unsupported analysis mode: {selected_mode}")
    if not files:
        raise HTTPException(status_code=400, detail="At least one document is required")

    extracted_sections: list[str] = []
    filenames: list[str] = []
    for upload in files:
        filename = upload.filename or "untitled"
        ext = ""
        if "." in filename:
            ext = "." + filename.rsplit(".", 1)[-1].lower()
        if ext not in _ALLOWED_EXTENSIONS:
            allowed = ", ".join(sorted(_ALLOWED_EXTENSIONS))
            raise HTTPException(status_code=400, detail=f"Unsupported file type: {ext}. Allowed: {allowed}")
        data = await upload.read()
        text = _extract_text_for_extension(ext, data).strip()
        if not text:
            continue
        filenames.append(filename)
        extracted_sections.append(f"File: {filename}\n{text[:32000]}")

    if not extracted_sections:
        raise HTTPException(status_code=400, detail="No readable text could be extracted from the uploaded files")

    label = title or ", ".join(filenames[:2]) or "Document set"
    mission_title = f"Document Intel · {label}"
    _update_document_mission(
        request,
        title=mission_title,
        summary=f"Preparing {selected_mode.replace('_', ' ')} analysis for {label}.",
        next_step="Wait for the document analysis to finish, then review the resulting brief.",
        phase="act",
    )

    api_key = os.environ.get("OPENAI_API_KEY", "").strip()
    if not api_key:
        _update_document_mission(
            request,
            title=mission_title,
            summary=f"Document analysis blocked for {label}.",
            next_step="Configure OPENAI_API_KEY, then retry the document analysis.",
            result="Document analysis requires OPENAI_API_KEY.",
            status="blocked",
            phase="retry",
            retry_hint="Add an OpenAI API key and rerun the document analysis.",
        )
        raise HTTPException(status_code=503, detail="Document analysis requires OPENAI_API_KEY")

    try:
        from openai import OpenAI
    except ImportError as exc:
        raise HTTPException(status_code=503, detail=f"openai package unavailable: {exc}") from exc

    model = os.environ.get("OPENJARVIS_DOCUMENT_MODEL", os.environ.get("OPENJARVIS_VISION_MODEL", "gpt-4o-mini"))
    prompt = _document_analysis_prompt(selected_mode, title or "", filenames)
    document_body = "\n\n---\n\n".join(extracted_sections)
    try:
        client = OpenAI(api_key=api_key)
        response = client.chat.completions.create(
            model=model,
            messages=[
                {"role": "system", "content": prompt},
                {
                    "role": "user",
                    "content": (
                        f"Analyze these documents for mode={selected_mode}.\n\n"
                        f"{document_body[:120000]}"
                    ),
                },
            ],
            temperature=0.2,
        )
    except Exception as exc:
        _update_document_mission(
            request,
            title=mission_title,
            summary=f"Document analysis failed for {label}.",
            next_step="Review the analysis error and retry with the same or a narrower document set.",
            result=str(exc),
            status="blocked",
            phase="retry",
            retry_hint="Retry after checking the model, API key, or uploaded documents.",
        )
        raise HTTPException(status_code=400, detail=f"Document analysis failed: {exc}") from exc

    content = ""
    choice = response.choices[0] if response.choices else None
    if choice and choice.message:
        content = (choice.message.content or "").strip()

    try:
        operator_memory = get_operator_memory_manager(request)
        operator_memory.add_explicit_memory(
            f"Document analysis ({selected_mode}) for {title or ', '.join(filenames[:3])}: {content[:600]}",
            tags=["document", selected_mode, "business-finance"],
        )
    except Exception:
        logger.debug("Document analysis memory update skipped", exc_info=True)

    _update_document_mission(
        request,
        title=mission_title,
        summary=f"Document analysis ready for {label}.",
        next_step="Review the document brief and turn it into a task, planner handoff, or executive decision.",
        result=content[:1000],
        status="complete",
        phase="done",
    )

    return DocumentAnalysisResponse(
        mode=selected_mode,
        content=content,
        files=filenames,
        model=model,
    )


@router.post("/export")
async def export_document_analysis(req: DocumentExportRequest):
    content = req.content.strip()
    if not content:
        raise HTTPException(status_code=400, detail="Export content is required")
    export_format = req.format.strip().lower()
    title = req.title.strip() or "Document Intel Export"
    stem = _safe_export_stem(title)
    if export_format == "docx":
        payload = _export_document_docx(title, req.mode, content)
        media_type = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        filename = f"{stem}.docx"
    elif export_format == "xlsx":
        payload = _export_document_xlsx(title, req.mode, content)
        media_type = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        filename = f"{stem}.xlsx"
    elif export_format == "txt":
        payload = _export_document_txt(title, req.mode, content)
        media_type = "text/plain; charset=utf-8"
        filename = f"{stem}.txt"
    else:
        raise HTTPException(status_code=400, detail=f"Unsupported export format: {export_format}")
    return StreamingResponse(
        io.BytesIO(payload),
        media_type=media_type,
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )
